#!/usr/bin/env python3
"""Build detailed PPT of Varsha chatbot changes (11 Aug 2026)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

OUT = "/Users/tndrra/Documents/IRAINS/iru-irains-backend/docs/Varsha_Chatbot_Changes_11_Aug_2026.pptx"

NAVY = RGBColor(0x00, 0x24, 0x67)
SKY = RGBColor(0x1A, 0x8F, 0xB8)
INK = RGBColor(0x0C, 0x1F, 0x3A)
MUTED = RGBColor(0x5A, 0x6F, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF9, 0xFC)
ACCENT = RGBColor(0x0A, 0x4A, 0x8A)


def set_run(run, size=14, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run(run, 22, True, WHITE)


def add_footer(slide, page, total):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"iRAINS · Varsha Chatbot · Changes 11 Aug 2026 · {page}/{total}"
    set_run(run, 10, False, MUTED)
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, left, top, width, height, items, size=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = item.get("level", 0)
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = item["text"]
        set_run(run, item.get("size", size), item.get("bold", False), item.get("color", INK))


def add_card(slide, left, top, width, height, title, body_lines, accent=SKY):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD0, 0xDC, 0xE8)
    # accent strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + 0.12), Inches(width - 0.4), Inches(0.35)
    )
    tr = title_box.text_frame.paragraphs[0].add_run()
    tr.text = title
    set_run(tr, 14, True, NAVY)

    body = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + 0.45), Inches(width - 0.4), Inches(height - 0.55)
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        set_run(run, 12, False, INK)


def title_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    # accent band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.6), Inches(13.333), Inches(1.9)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Varsha Rainfall Chatbot"
    set_run(r, 36, True, WHITE)

    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.8))
    r2 = s.text_frame.paragraphs[0].add_run()
    r2.text = "Detailed Change Summary — Frontend & Backend"
    set_run(r2, 22, False, RGBColor(0xD4, 0xE8, 0xF5))

    m = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.2))
    tf = m.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Session date: 11 August 2026  ·  Product: iRAINS / IMD"
    set_run(run, 14, False, WHITE)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Repos: iru-irains-backend + IRAINS-frontend  ·  Feature: Clarification flows, category answers, UI polish"
    set_run(run2, 13, False, RGBColor(0xD4, 0xE8, 0xF5))


def agenda_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Agenda")
    add_bullets(
        slide,
        0.7,
        1.3,
        11.5,
        5.5,
        [
            {"text": "1. Overview — what we built yesterday"},
            {"text": "2. Clarification flow (Did you mean → Period → Data)"},
            {"text": "3. Backend changes (clarification.js, chatService, apiExecutor)"},
            {"text": "4. Frontend changes (Varsha UI, chips, answer formatting)"},
            {"text": "5. Bug fixes from live testing (with before/after)"},
            {"text": "6. Category miss messaging (Excess vs actual category)"},
            {"text": "7. Date / period handling (June month, Historical, Season so far)"},
            {"text": "8. Files touched & test plan"},
            {"text": "9. Summary & next recommendations"},
        ],
        size=18,
    )
    add_footer(slide, 2, total)


def overview_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "1. Overview — Yesterday’s Goals")
    add_card(
        slide,
        0.5,
        1.3,
        6.0,
        5.3,
        "Product intent",
        [
            "Make Varsha behave like a guided rainfall assistant, not a free-form LLM.",
            "",
            "• Catch typos before calling APIs",
            "• Ask Yes/No for place corrections",
            "• Ask for period when timeframe is unclear",
            "• Answer category questions correctly",
            "• Never contradict itself (e.g. show departure then say “no data”)",
            "• Polish clarify UI (chips, hints, no raw HTML)",
        ],
        SKY,
    )
    add_card(
        slide,
        6.8,
        1.3,
        5.9,
        5.3,
        "Outcome",
        [
            "New clarification layer on backend",
            "Frontend clarify chips + free-text matching",
            "Deterministic answers for category misses",
            "Safer location fuzzy-matching",
            "Correct month/year & historical ranges",
            "Professional clarify message layout",
            "",
            "Backend ~+578 LOC  |  Frontend ~+837 LOC",
        ],
        ACCENT,
    )
    add_footer(slide, 3, total)


def flow_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "2. Target User Flow")
    add_bullets(
        slide,
        0.6,
        1.2,
        12,
        1.0,
        [
            {
                "text": "Example: “can you give me chenai data of excess” / “no rain” / “large excess”",
                "bold": True,
                "size": 15,
            }
        ],
    )

    steps = [
        ("1", "User question", "Typo + category\n+ unclear period"),
        ("2", "Did you mean?", "Yes / No chips\nfor Chennai"),
        ("3", "Pick period", "Today / Yesterday\nThis month / …"),
        ("4", "Optional month", "Specific month\nor type March 2023"),
        ("5", "Data answer", "Category hit\nor category miss"),
    ]
    x = 0.45
    for num, title, body in steps:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(2.3), Inches(2.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(0xC5, 0xD5, 0xE5)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.85), Inches(2.55), Inches(0.55), Inches(0.55)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        br = badge.text_frame.paragraphs[0]
        br.alignment = PP_ALIGN.CENTER
        run = br.add_run()
        run.text = num
        set_run(run, 14, True, WHITE)

        tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.25), Inches(2.1), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        set_run(r, 13, True, NAVY)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        set_run(r2, 11, False, MUTED)
        x += 2.55

    note = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12), Inches(1.2))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Rule: never skip Did-you-mean when a place looks misspelled. Never invent a place from English words (e.g. “than” ≠ Thane). Period chips must not glue onto the place name (e.g. “goa monthly”)."
    set_run(r, 13, False, INK)
    add_footer(slide, 4, total)


def backend_overview(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "3. Backend — Files & Responsibilities")
    add_card(
        slide,
        0.45,
        1.25,
        4.1,
        5.4,
        "NEW  clarification.js",
        [
            "Pre-chat validation layer",
            "• Fuzzy place matching",
            "• Did-you-mean Yes/No",
            "• Invalid date / location",
            "• Suspicious mm values",
            "• Mixed temp+rain intent",
            "• Category period chips",
            "• Specific-month picker",
            "• Location master (DB + seed)",
            "• healPlaceFilterLevel()",
        ],
        SKY,
    )
    add_card(
        slide,
        4.7,
        1.25,
        4.1,
        5.4,
        "UPDATED  chatService.js",
        [
            "Wires clarification into chat",
            "• runPreChatClarifications()",
            "• Category-miss answer rules",
            "• Override LLM when empty",
            "  category filter (no “data",
            "  not available” contradiction)",
            "• Prompt: category_miss ≠",
            "  missing rainfall data",
            "• Strip confirmed_rainfall_mm",
            "  token before planning",
        ],
        ACCENT,
    )
    add_card(
        slide,
        8.95,
        1.25,
        3.9,
        5.4,
        "UPDATED  apiExecutor.js",
        [
            "Date & category healers",
            "• extractCategoriesFromQuestion",
            "• sanitizeDatesFromQuestion",
            "• applyMonthRangeFromQuestion",
            "• Historical → season start",
            "• category_miss payload",
            "• State/district fallbacks",
            "• sanitizeRainfallAction",
            "",
            "Docs: IRAINS_API_CATALOG.md",
            "(Q5c place + category)",
        ],
        RGBColor(0x2A, 0x6F, 0xA8),
    )
    add_footer(slide, 5, total)


def backend_clarify_detail(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "3a. Backend Clarification Types")
    rows = [
        ("did_you_mean", "chenai → Chennai", "Yes / No chips; Yes rewrites place then continues"),
        ("ambiguous_timeframe", "place known, no period", "Today, Yesterday, Last 7 days, This/Last month, Specific month / Season so far"),
        ("which_month", "Specific month picked", "Last 12 months + Year 20XX chips; free text e.g. March 2023"),
        ("invalid_location", "unknown place", "Suggest nearby district names as chips"),
        ("invalid_date", "31 Feb etc.", "Ask for a valid date"),
        ("suspicious_rainfall_value", "4000mm", "Yes 40 mm / Keep 4000 mm; no re-ask loop"),
        ("mixed_concept", "temp + rain", "Rainfall available; Temperature disabled"),
    ]
    y = 1.2
    for typ, trigger, behavior in rows:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(y), Inches(12.4), Inches(0.72)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xD0, 0xDC, 0xE8)
        t = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.08), Inches(3.2), Inches(0.55))
        r = t.text_frame.paragraphs[0].add_run()
        r.text = typ
        set_run(r, 12, True, NAVY)
        t2 = slide.shapes.add_textbox(Inches(3.9), Inches(y + 0.08), Inches(3.0), Inches(0.55))
        r2 = t2.text_frame.paragraphs[0].add_run()
        r2.text = trigger
        set_run(r2, 12, False, SKY)
        t3 = slide.shapes.add_textbox(Inches(7.0), Inches(y + 0.08), Inches(5.6), Inches(0.55))
        tf = t3.text_frame
        tf.word_wrap = True
        r3 = tf.paragraphs[0].add_run()
        r3.text = behavior
        set_run(r3, 11, False, INK)
        y += 0.78
    add_footer(slide, 6, total)


def backend_location(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "3b. Location Matching — Safeguards")
    add_bullets(
        slide,
        0.6,
        1.25,
        12,
        5.5,
        [
            {"text": "Problems fixed", "bold": True, "size": 16, "color": NAVY},
            {"text": "• “of no rain” was parsed as place “No” → skipped Chennai typo check"},
            {"text": "• “more than 40mm” matched “than” → “Did you mean Thane?”"},
            {"text": "• “rainfall in goa” + Monthly became invalid location “goa monthly”"},
            {"text": "", "size": 8},
            {"text": "Fixes implemented", "bold": True, "size": 16, "color": NAVY},
            {"text": "• Strip category phrases (no rain / large excess / …) before place extract"},
            {"text": "• Prefer “give me X data” pattern over loose “of …” prep capture"},
            {"text": "• Expanded STOP_TOKENS (than, more, any, has, monthly, historical, …)"},
            {"text": "• Tightened substring similarity (than ⊂ thane no longer scores 0.92)"},
            {"text": "• Skip typo scan on threshold listing questions without a place cue"},
            {"text": "• Always run typo scan even if a bad placeHint was captured"},
            {"text": "• Period words stripped from extracted place names"},
        ],
        size=14,
    )
    add_footer(slide, 7, total)


def backend_dates(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "3c. Date & Period Handling")
    add_card(
        slide,
        0.45,
        1.25,
        6.1,
        5.4,
        "sanitizeDatesFromQuestion / applyMonthRange",
        [
            "Relative phrases always win over LLM-copied catalog dates.",
            "",
            "• today / yesterday → single day",
            "• last 7 days / this week → rolling week",
            "• this month / monthly → 1st → today",
            "• last month → full previous month",
            "• historical / season so far / seasonal",
            "  → SW monsoon start (1 June) → today",
            "",
            "Whole-month phrases (current year if omitted):",
            "• in June / on June / June month / month of June",
            "• June 2024 keeps explicit year",
            "• Prevents LLM inventing 2023-06-01…06-30",
        ],
        SKY,
    )
    add_card(
        slide,
        6.8,
        1.25,
        5.9,
        5.4,
        "Live bug: Goa June",
        [
            "User: rainfall in goa on june month",
            "",
            "Before:",
            "• Answered June 2023 (catalog example)",
            "• All values “Not available”",
            "",
            "After:",
            "• Forces 2026-06-01 → 2026-06-30",
            "  (server year when year omitted)",
            "",
            "Also renamed vague “Historical” chip",
            "→ “Season so far” so the range is clear.",
        ],
        ACCENT,
    )
    add_footer(slide, 8, total)


def backend_category(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "3d. Category Questions & category_miss")
    add_bullets(
        slide,
        0.6,
        1.2,
        12,
        5.6,
        [
            {"text": "Supported categories", "bold": True, "size": 15, "color": NAVY},
            {"text": "Large Excess · Excess · Deficient · Large Deficient · No Rain (+ Normal with context)"},
            {"text": "Tolerates typos like excesss; heals category from question text over LLM drift.", "size": 13},
            {"text": "", "size": 6},
            {"text": "When asked category but place was different", "bold": True, "size": 15, "color": NAVY},
            {"text": "API returns empty after filter_by_departure_category, but unfiltered row exists → category_miss { category, departure, wanted }."},
            {"text": "", "size": 6},
            {"text": "Deterministic answer (LLM overridden)", "bold": True, "size": 15, "color": NAVY},
            {"text": "“Chennai was not in Excess for 2026-06-01 to 2026-06-30. It was Deficient (departure -25.1%). Rainfall data is available — it just was not Excess.”"},
            {"text": "", "size": 6},
            {"text": "Prompt rule", "bold": True, "size": 15, "color": NAVY},
            {"text": "If category_miss is present → NEVER say “data is not available”. Only say that when sample is empty AND category_miss is absent."},
        ],
        size=14,
    )
    add_footer(slide, 9, total)


def frontend_overview(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "4. Frontend — Files Touched")
    add_card(
        slide,
        0.45,
        1.25,
        6.1,
        5.4,
        "rainfall-chatbot.component.ts / html / css",
        [
            "• PendingBackendClarify state machine",
            "• handleClarifyResponse + formatClarifyAnswer",
            "• Chip variants: month / year / action / primary",
            "• Free-text matching while chips are open",
            "  (Yes/No, periods, Yes 40 mm, months)",
            "• buildClarifiedQuestion rewrites:",
            "  – did_you_mean place replace",
            "  – period / month append (safe)",
            "  – suspicious mm (handles 4000mm)",
            "• Category-empty card from category_miss",
            "• Source-page link (“Data from: …”)",
        ],
        SKY,
    )
    add_card(
        slide,
        6.8,
        1.25,
        5.9,
        5.4,
        "rainfall-chat.service.ts + env",
        [
            "• Types for clarify payloads",
            "  (did_you_mean, options, from/to,",
            "   original_value, category_miss)",
            "• Ask API contract stays POST",
            "  /ollama-chat (or project route)",
            "• Environment tip for local API base",
            "",
            "UI polish",
            "• Title + muted hint (no raw <em>)",
            "• No pipe-list of options in bubble",
            "• Chips separated by top border",
            "• Year chips dashed; Yes primary",
        ],
        ACCENT,
    )
    add_footer(slide, 10, total)


def frontend_ui(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "4a. Clarify UI — Before vs After")
    add_card(
        slide,
        0.45,
        1.25,
        6.1,
        5.4,
        "Before (unprofessional)",
        [
            "• Raw HTML visible: <em>March 2023</em>",
            "• Options dumped as pipe text:",
            "  August 2026 | July 2026 | …",
            "• Flat same-weight paragraph",
            "• Pill chips only, weak hierarchy",
            "• Confusing contradictory copy",
            "  (“It was Deficient…” +",
            "   “data is not available”)",
        ],
        RGBColor(0x8A, 0x3A, 0x3A),
    )
    add_card(
        slide,
        6.8,
        1.25,
        5.9,
        5.4,
        "After (professional)",
        [
            "• clarify-title + clarify-hint",
            "• Options only as interactive chips",
            "• Month chips soft cards",
            "• Year chips dashed secondary",
            "• Yes = primary navy button",
            "• Category miss: clear 3-line card",
            "  1) not in requested category",
            "  2) actual category + departure",
            "  3) data exists — wrong category",
        ],
        RGBColor(0x1F, 0x7A, 0x4D),
    )
    add_footer(slide, 11, total)


def frontend_rewrite(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "4b. Frontend Question Rewrites")
    add_bullets(
        slide,
        0.6,
        1.2,
        12,
        5.6,
        [
            {"text": "did_you_mean + Yes", "bold": True, "size": 15, "color": NAVY},
            {"text": "Replace typo token: “chenai …” → “Chennai …” then re-ask backend (period chips next)."},
            {"text": "", "size": 6},
            {"text": "ambiguous_timeframe / which_month", "bold": True, "size": 15, "color": NAVY},
            {"text": "Aliases: monthly→this month, historical→season so far."},
            {"text": "Insert period into rainfall phrase: “this month rainfall in goa” (avoids “goa monthly”)."},
            {"text": "Specific month → “month of March 2024”; Year chip → “specific_month 2023”."},
            {"text": "", "size": 6},
            {"text": "suspicious_rainfall_value", "bold": True, "size": 15, "color": NAVY},
            {"text": "Rewrite 4000mm / 4000 mm → 40 mm. Keep 4000 adds confirmed_rainfall_mm so backend won’t loop."},
            {"text": "", "size": 6},
            {"text": "Free-text while chips open", "bold": True, "size": 15, "color": NAVY},
            {"text": "Matches Yes/No, period labels, “Yes, 40 mm”, typed months — otherwise nudges user to tap a chip."},
        ],
        size=14,
    )
    add_footer(slide, 12, total)


def bugs_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "5. Live Bugs Fixed Yesterday")
    bugs = [
        ("chenai + no rain → “for No?”", "Skipped Did-you-mean; place=No from “no rain”", "Typo first; scrub categories from place extract"),
        ("Yes 40 mm looped on 4000mm", "\\b4000\\b failed on glued 4000mm", "Replace NNmm unit; confirmed_rainfall_mm for Keep"),
        ("“more than 40mm” → Thane?", "than fuzzy-matched Thane", "Stop words + stricter similarity + skip threshold typos"),
        ("Goa + Monthly → goa monthly", "Period glued to place → invalid_location", "Period aliases + strip period from place; safer rewrite"),
        ("Historical = today only", "historical not in date sanitizer", "Map to season start→today; chip “Season so far”"),
        ("June month → June 2023", "LLM copied catalog year; “on june month” missed", "applyMonthRange recognizes on/june month; force current year"),
        ("Excess miss + “no data”", "LLM appended contradictory line", "Deterministic category_miss answer; override LLM"),
    ]
    y = 1.15
    headers = ("Issue", "Root cause", "Fix")
    for i, h in enumerate(headers):
        left = 0.45 + i * 4.2
        t = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(4.0), Inches(0.35))
        r = t.text_frame.paragraphs[0].add_run()
        r.text = h
        set_run(r, 12, True, MUTED)
    y = 1.5
    for issue, cause, fix in bugs:
        for i, text in enumerate((issue, cause, fix)):
            left = 0.45 + i * 4.2
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(4.05), Inches(0.7)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = RGBColor(0xD0, 0xDC, 0xE8)
            tb = slide.shapes.add_textbox(Inches(left + 0.1), Inches(y + 0.08), Inches(3.85), Inches(0.55))
            tf = tb.text_frame
            tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = text
            set_run(r, 10, i == 0, NAVY if i == 0 else INK)
        y += 0.75
    add_footer(slide, 13, total)


def files_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "6. Files Changed")
    add_card(
        slide,
        0.45,
        1.25,
        6.1,
        5.4,
        "Backend (iru-irains-backend)",
        [
            "NEW",
            "  controllers/ollamaChat/clarification.js",
            "",
            "MODIFIED",
            "  controllers/ollamaChat/chatService.js",
            "  controllers/ollamaChat/apiExecutor.js",
            "  docs/IRAINS_API_CATALOG.md",
            "",
            "Approx. +578 / −34 lines (tracked)",
        ],
        SKY,
    )
    add_card(
        slide,
        6.8,
        1.25,
        5.9,
        5.4,
        "Frontend (IRAINS-frontend)",
        [
            "MODIFIED",
            "  shared/rainfall-chatbot/",
            "    rainfall-chatbot.component.ts",
            "    rainfall-chatbot.component.html",
            "    rainfall-chatbot.component.css",
            "  services/rainfall-chat/",
            "    rainfall-chat.service.ts",
            "  environment/environment.ts",
            "",
            "Approx. +837 / −70 lines (tracked)",
        ],
        ACCENT,
    )
    add_footer(slide, 14, total)


def test_plan(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "7. Recommended Test Plan")
    add_bullets(
        slide,
        0.6,
        1.2,
        12,
        5.6,
        [
            {"text": "1. chenai + excess / no rain → Did you mean Chennai? → Yes → period → answer"},
            {"text": "2. Did you mean → No → asks to retype place"},
            {"text": "3. Specific month → chips + type “March 2023” → no month re-loop"},
            {"text": "4. Chennai Excess for June 2026 → category_miss copy (not “no data”) if Deficient"},
            {"text": "5. “more than 40mm” → lists districts; never “Did you mean Thane?”"},
            {"text": "6. “more than 4000mm” → Yes 40 mm continues once; Keep 4000 continues once"},
            {"text": "7. rainfall in goa → period chips → This month / Season so far → real range"},
            {"text": "8. rainfall in goa on june month → 2026-06-01..30 (not 2023)"},
            {"text": "9. Clarify UI shows title/hint + chips only (no <em>, no pipe list)"},
        ],
        size=15,
    )
    add_footer(slide, 15, total)


def summary_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "8. Summary")
    add_bullets(
        slide,
        0.6,
        1.25,
        12,
        5.5,
        [
            {"text": "Yesterday turned Varsha from “LLM answers” into a guided rainfall workflow.", "bold": True, "size": 16},
            {"text": ""},
            {"text": "Backend: clarification layer + date/category healers + deterministic category_miss answers."},
            {"text": "Frontend: chip UX, free-text matching, safe question rewrites, professional clarify formatting."},
            {"text": "Quality: fixed seven high-visibility live bugs (place typos, units, Thane false positive, Goa monthly, Historical, June year, contradictory no-data)."},
            {"text": ""},
            {"text": "Suggested next steps", "bold": True, "size": 15, "color": NAVY},
            {"text": "• Commit + PR for backend clarification.js and frontend chatbot changes"},
            {"text": "• Add automated unit tests for extractMentionedLocation / applyMonthRange / buildClarifiedQuestion"},
            {"text": "• Align all period chip sets (category vs place) to the same option list"},
            {"text": "• Consider logging clarify types for analytics"},
        ],
        size=14,
    )
    add_footer(slide, 16, total)


def end_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    t = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Thank you"
    set_run(r, 40, True, WHITE)
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11.3), Inches(1.2))
    p = s.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r2 = p.add_run()
    r2.text = "Varsha · iRAINS Rainfall Companion\nDetailed change pack — 11 August 2026"
    set_run(r2, 16, False, RGBColor(0xD4, 0xE8, 0xF5))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 17

    title_slide(prs, total)
    agenda_slide(prs, total)
    overview_slide(prs, total)
    flow_slide(prs, total)
    backend_overview(prs, total)
    backend_clarify_detail(prs, total)
    backend_location(prs, total)
    backend_dates(prs, total)
    backend_category(prs, total)
    frontend_overview(prs, total)
    frontend_ui(prs, total)
    frontend_rewrite(prs, total)
    bugs_slide(prs, total)
    files_slide(prs, total)
    test_plan(prs, total)
    summary_slide(prs, total)
    end_slide(prs, total)

    import os

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
